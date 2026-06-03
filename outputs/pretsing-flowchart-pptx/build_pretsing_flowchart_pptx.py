from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT = "/Users/vivekgupta/workspace/DYNAMIC/outputs/pretsing-flowchart-pptx/dynamic_flowchart_A4_editable.pptx"


prs = Presentation()
prs.slide_width = Inches(8.27)
prs.slide_height = Inches(11.69)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def add_box(name, x, y, w, h, text, size=8.2):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.8)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLACK
    return shape


def add_label(name, x, y, w, h, text, size=6.6, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLACK
    return shape


def arrow(name, x1, y1, x2, y2, width=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.name = name
    c.shadow.inherit = False
    c.line.color.rgb = BLACK
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def line(name, x1, y1, x2, y2, width=1.4, dash=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.name = name
    c.shadow.inherit = False
    c.line.color.rgb = BLACK
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def elbow_down_right(name, x, y, mid_y, end_x, end_y):
    line(name + "_v", x, y, x, mid_y)
    line(name + "_h", x, mid_y, end_x, mid_y)
    arrow(name + "_end", end_x, mid_y, end_x, end_y)


# Main nodes, placed to match the source flowchart on a portrait A4 slide.
hhq = add_box("Baseline Household Questionnaire", 3.55, 1.78, 0.85, 0.62, "Baseline\nHousehold\nQuestionnaire")
wq = add_box("Baseline Woman's Questionnaire", 3.55, 2.70, 0.85, 0.62, "Baseline\nWoman's\nQuestionnaire")
pef = add_box("Pregnancy Enrollment Form", 3.55, 3.63, 0.85, 0.62, "Pregnancy\nEnrollment\nForm")
pff1 = add_box("Pregnancy Follow-Up Form 1", 3.55, 4.55, 0.85, 0.62, "Pregnancy\nFollow-Up\nForm 1")
pffn = add_box("Pregnancy Follow-Up Form N", 3.55, 5.45, 0.85, 0.62, "Pregnancy\nFollow-Up\nForm N")
pof = add_box("Pregnancy Outcome Form", 3.55, 6.45, 0.85, 0.62, "Pregnancy\nOutcome\nForm")
baf = add_box("Birth Assessment Form", 3.55, 7.35, 0.85, 0.62, "Birth\nAssessment\nForm")
nff1 = add_box("Child Follow-Up Form 1", 3.55, 8.25, 0.85, 0.62, "Child\nFollow-Up\nForm 1")
nffn = add_box("Child Follow-Up Form N", 3.55, 9.35, 0.85, 0.62, "Child\nFollow-Up\nForm N")
end_study = add_box("End of Study Period", 3.34, 10.35, 1.28, 0.34, "END OF STUDY PERIOD", 7.0)

hrf = add_box("Household Rounds Form", 1.12, 2.70, 0.90, 0.62, "Household\nRounds\nForm 1,...,N")
end_left = add_box("End of Study Period Left", 0.90, 4.67, 1.25, 0.36, "END OF STUDY PERIOD", 7.0)
uf = add_box("Ultrasound Form", 6.18, 4.55, 0.88, 0.62, "Ultrasound\nForm")
end_abort = add_box("End Abort Fetal Death", 6.18, 6.85, 0.88, 0.62, "END")
sbf = add_box("Stillbirth Form", 6.18, 7.73, 0.88, 0.62, "Stillbirth\nForm")
cdf = add_box("Child Death Form", 6.18, 8.72, 0.88, 0.62, "Child\nDeath\nForm")
va = add_box("Verbal Autopsy", 6.18, 9.62, 0.88, 0.62, "VA")

# Labels from source image.
add_label("No Eligible Women Label", 1.35, 1.62, 1.50, 0.34, "Household without eligible women\n(= women at risk of becoming pregnant\nduring study period)", 6.0)
add_label("Ever Married Label", 4.92, 2.42, 1.55, 0.16, "Ever-married woman aged 18-49", 6.0, PP_ALIGN.LEFT)
add_label("Existing Eligible Woman Label", 2.33, 2.75, 1.15, 0.14, "Existing eligible woman", 6.0)
add_label("New Eligible Woman Label", 2.38, 3.22, 1.00, 0.14, "New eligible woman", 6.0)
add_label("Pregnancy Detected Right Label", 4.92, 3.33, 1.05, 0.14, "Pregnancy detected", 6.0, PP_ALIGN.LEFT)
add_label("Pregnancy Detected Left Label", 2.16, 3.90, 1.18, 0.14, "Pregnancy detected", 6.0)
add_label("USG Enrollment Label", 5.22, 4.05, 0.88, 0.34, "1st USG report\navailable", 6.0)
add_label("USG FU1 Label", 4.68, 4.58, 1.04, 0.22, "1st USG report\navailable", 6.0)
add_label("USG FUN Label", 5.22, 5.98, 0.88, 0.34, "1st USG report\navailable", 6.0)
add_label("Delivery Occurs Label", 2.25, 6.96, 1.00, 0.16, "Delivery Occurs", 6.0)
add_label("Abortion Label", 5.05, 6.75, 1.00, 0.24, "Abortion or\nfetal death < 20w", 6.0)
add_label("Fetal Death Label", 5.20, 7.74, 0.95, 0.14, "Fetal death >= 20w", 6.0)
add_label("Live Birth Death Label", 5.22, 8.30, 0.88, 0.24, "Death of a live\nbirth", 6.0)
add_label("Child Survived Label", 2.25, 8.45, 1.00, 0.16, "Child Survived", 6.0)
add_label("Child Died FU1 Label", 5.00, 8.80, 1.10, 0.14, "Child died", 6.0)
add_label("Child Died FUN Label", 5.00, 9.92, 1.10, 0.14, "Child died", 6.0)

# Vertical core flow.
arrow("HHQ to WQ", 3.98, 2.40, 3.98, 2.70)
arrow("WQ to PEF", 3.98, 3.32, 3.98, 3.63)
arrow("PEF to PFF1", 3.98, 4.25, 3.98, 4.55)
line("PFF1 to PFFN dotted", 3.98, 5.17, 3.98, 5.45, dash=True)
arrow("PFFN to POF", 3.98, 6.07, 3.98, 6.45)
arrow("POF to BAF", 3.98, 7.07, 3.98, 7.35)
arrow("BAF to NFF1", 3.98, 7.97, 3.98, 8.25)
line("NFF1 to NFFN dotted", 3.98, 8.87, 3.98, 9.35, dash=True)
arrow("NFFN to Study End", 3.98, 9.97, 3.98, 10.35)

# Household rounds and baseline branches.
line("NoEligible down", 1.50, 1.96, 1.50, 2.18)
line("NoEligible across", 1.50, 2.18, 3.55, 2.18)
arrow("NoEligible to HRF", 1.50, 2.18, 1.50, 2.70)
arrow("HHQ Existing to HRF", 3.55, 2.85, 2.02, 2.85)
arrow("HRF New to WQ", 2.02, 3.16, 3.55, 3.16)
arrow("HRF Pregnancy to PEF", 2.02, 3.32, 3.55, 3.94)
arrow("HRF to End", 1.50, 3.32, 1.50, 4.67)
arrow("Ever Married to WQ", 5.02, 2.55, 4.40, 2.55)
arrow("Preg Detected to PEF", 5.02, 3.48, 4.40, 3.48)

# Ultrasound branches.
arrow("PEF to UF", 4.40, 3.94, 6.18, 4.86)
arrow("PFF1 to UF", 4.40, 4.86, 6.18, 4.86)
arrow("PFFN to UF", 4.40, 5.76, 6.18, 4.86)

# Delivery/outcome branches.
arrow("Delivery to POF", 3.16, 6.91, 3.55, 6.91)
arrow("BAF to Abort End", 4.40, 7.66, 6.18, 7.16)
arrow("BAF to SBF", 4.40, 7.76, 6.18, 8.04)
arrow("BAF to CDF", 4.40, 7.86, 6.18, 9.03)
arrow("Child Survived to NFF1", 3.16, 8.56, 3.55, 8.56)
arrow("NFF1 to CDF", 4.40, 8.56, 6.18, 9.03)
arrow("NFFN to CDF", 4.40, 9.66, 6.18, 9.03)
arrow("CDF to VA", 6.62, 9.34, 6.62, 9.62)

# Stillbirth to VA elbow at right side.
line("SBF right", 7.06, 8.04, 7.42, 8.04)
line("SBF down", 7.42, 8.04, 7.42, 9.93)
arrow("SBF to VA", 7.42, 9.93, 7.06, 9.93)

# Add unobtrusive slide note marker via authoring properties.
prs.core_properties.title = "DYNAMIC Flowchart A4 Editable"
prs.core_properties.subject = "Editable PowerPoint recreation of dynamic_flowchart_v2.pdf"

prs.save(OUT)
print(OUT)
